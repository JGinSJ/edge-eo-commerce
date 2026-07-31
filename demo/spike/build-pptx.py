#!/usr/bin/env python3
"""Build the customer-facing prioritized-delivery deck as a native .pptx.

The HTML deck (customer-deck.html) is the source of truth for content; this
produces the same twelve slides as real PowerPoint objects — text frames, shapes
and a native table — so the deck can be edited in PowerPoint or uploaded to
Google Slides (which imports .pptx directly). Exporting the HTML to PDF would be
one keystroke, but a PDF cannot be edited by the people who have to present it.

Figures come from `POST /spike/skinnify` against the telco support fixture and
are duplicated here as constants; if the transforms change, re-measure and
update BOTH files.

    pip install python-pptx && python3 build-pptx.py [out.pptx]
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Palette: the HTML deck's light-theme tokens ──────────────────────────────
GROUND       = RGBColor(0xF6, 0xF8, 0xFA)
PANEL        = RGBColor(0xFF, 0xFF, 0xFF)
INK          = RGBColor(0x0A, 0x1A, 0x2F)
INK2         = RGBColor(0x33, 0x4A, 0x61)
MUTE         = RGBColor(0x7B, 0x8B, 0x9A)
LINE         = RGBColor(0xDC, 0xE4, 0xEB)
EDGE         = RGBColor(0x00, 0x89, 0xC7)
EDGE_SOFT    = RGBColor(0xE3, 0xF2, 0xFB)
BURIED       = RGBColor(0xC9, 0x6A, 0x00)
REACHED      = RGBColor(0x0E, 0x8C, 0x61)
REACHED_SOFT = RGBColor(0xE1, 0xF4, 0xEC)
RAIL         = RGBColor(0xE7, 0xED, 0xF2)

# Arial and Consolas both ship with Office on Windows and macOS; Google Slides
# substitutes cleanly. Deliberately not the CSS stack — system-ui means nothing here.
F_DISPLAY = "Arial"
F_BODY    = "Arial"
F_MONO    = "Consolas"

SLIDE_W, SLIDE_H = 13.333, 7.5
M = 0.75                      # side margin
CW = SLIDE_W - 2 * M          # content width


# ── Primitives ───────────────────────────────────────────────────────────────
def textbox(slide, x, y, w, h, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tf


def write(tf, text, *, size=14, bold=False, color=INK, font=F_BODY,
          spacing=None, caps=False, line=None, space_after=0, para=None,
          align=None):
    """Append a run to `para`, or to the frame's first/next paragraph."""
    p = para if para is not None else (
        tf.paragraphs[0] if not tf.paragraphs[0].runs and len(tf.paragraphs) == 1
        else tf.add_paragraph())
    if align is not None:
        p.alignment = align          # added paragraphs do NOT inherit the frame's
    r = p.add_run()
    r.text = text.upper() if caps else text
    f = r.font
    f.size, f.bold, f.name = Pt(size), bold, font
    f.color.rgb = color
    if spacing is not None:                      # letter-spacing, in points
        r.font._rPr.set('spc', str(int(spacing * 100)))
    if line:
        p.line_spacing = line
    p.space_after = Pt(space_after)
    return p


def rect(slide, x, y, w, h, *, fill=None, border=None, width=1.0,
         rounded=False, dash=False, alpha=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        shp.adjustments[0] = 0.06
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if alpha is not None:
            # python-pptx has no transparency API. Ruler bars need it: at full
            # opacity the answer marker disappears into a bar of its own colour.
            from pptx.oxml.ns import qn
            from lxml import etree
            clr = shp.fill._xPr.find(qn('a:solidFill')).find(qn('a:srgbClr'))
            etree.SubElement(clr, qn('a:alpha')).set('val', str(int(alpha * 1000)))
    if border is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = border
        shp.line.width = Pt(width)
        if dash:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return shp


def arrow(slide, x1, y1, x2, y2, *, color=MUTE, width=1.5, dash=False):
    """Straight connector with an arrowhead, drawn as a thin shape."""
    from pptx.enum.shapes import MSO_CONNECTOR
    cx = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cx.line.color.rgb = color
    cx.line.width = Pt(width)
    if dash:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        cx.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    # python-pptx has no arrowhead API; set it on the underlying line element.
    ln = cx.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    from lxml import etree
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('h', 'med')
    return cx


# ── Slide furniture ──────────────────────────────────────────────────────────
def new_slide(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    bg = s.background.fill
    bg.solid(); bg.fore_color.rgb = GROUND
    # progress rail along the very top — the HTML deck's brand bar
    rect(s, 0, 0, SLIDE_W * (n / total), 0.045, fill=EDGE)
    tf = textbox(s, M, SLIDE_H - 0.5, 2.0, 0.28)
    write(tf, "%02d / %02d" % (n, total), size=9, color=MUTE, font=F_MONO, spacing=1.2)
    return s


def head(s, eyebrow, title, *, lede=None, size=30, y=0.62):
    tf = textbox(s, M, y, CW, 0.3)
    write(tf, eyebrow, size=9, bold=True, color=EDGE, font=F_MONO, spacing=1.6, caps=True)
    tf2 = textbox(s, M, y + 0.4, CW * 0.86, 1.0)
    write(tf2, title, size=size, bold=True, color=INK, font=F_DISPLAY, line=1.04)
    # Titles set their own breaks, but estimate wrapping too — a title that wraps
    # unexpectedly used to run straight through the lede beneath it. ~48 characters
    # per line at 30pt bold across 0.86 of the content width.
    per_line = 48 if size >= 30 else 56
    lines = sum(max(1, -(-len(seg) // per_line)) for seg in title.split("\n"))
    bottom = y + 0.4 + (0.52 if size >= 30 else 0.44) * lines
    if lede:
        tf3 = textbox(s, M, bottom + 0.42, CW * 0.72, 0.9)
        write(tf3, lede, size=13, color=INK2, font=F_BODY, line=1.45)
        bottom += 1.0
    return bottom


def card(s, x, y, w, h, *, tag=None, title=None, body=None, accent=None,
         bullets=None):
    shp = rect(s, x, y, w, h, fill=PANEL, border=accent or LINE,
               width=1.5 if accent else 1.0, rounded=True)
    pad = 0.22
    ty = y + pad
    if tag:
        tf = textbox(s, x + pad, ty, w - 2 * pad, 0.22)
        write(tf, tag, size=8, bold=True, color=accent or MUTE, font=F_MONO,
              spacing=1.2, caps=True)
        ty += 0.34
    if title:
        tf = textbox(s, x + pad, ty, w - 2 * pad, 0.3)
        write(tf, title, size=13, bold=True, color=INK, font=F_DISPLAY, line=1.1)
        ty += 0.30 + 0.20 * title.count("\n")
    if body:
        tf = textbox(s, x + pad, ty + 0.06, w - 2 * pad, h - (ty - y) - pad)
        write(tf, body, size=10.5, color=INK2, font=F_BODY, line=1.4)
    if bullets:
        tf = textbox(s, x + pad, ty + 0.06, w - 2 * pad, h - (ty - y) - pad)
        for i, b in enumerate(bullets):
            p = None if i == 0 else tf.add_paragraph()
            write(tf, "▪  " + b, size=10.5, color=INK2, font=F_BODY,
                  line=1.35, space_after=6, para=p)
    return shp


def note(s, text, y):
    tf = textbox(s, M, y, CW * 0.8, 0.5)
    write(tf, text, size=9.5, color=MUTE, font=F_BODY, line=1.4)


# ── The token ruler ──────────────────────────────────────────────────────────
# bar length = payload, marker = where the answer sits. Both are the measurement.
def ruler(s, y, rows, *, track_x=2.35, track_w=8.1, row_h=0.46, gap=0.30):
    for i, (name, bar_pct, mark_pct, total, ans, tone) in enumerate(rows):
        ry = y + i * (row_h + gap)
        tf = textbox(s, M, ry + row_h / 2 - 0.11, track_x - M - 0.2, 0.24)
        write(tf, name, size=9, bold=True, color=INK2, font=F_MONO, spacing=1.0, caps=True)
        rect(s, track_x, ry, track_w, row_h, fill=RAIL)
        rect(s, track_x, ry, track_w * bar_pct, row_h, fill=tone, alpha=42)
        rect(s, track_x + track_w * mark_pct - 0.015, ry - 0.07, 0.03, row_h + 0.14,
             fill=INK if tone == EDGE else tone)
        tf = textbox(s, track_x + track_w + 0.25, ry - 0.02, 1.9, 0.5, align=PP_ALIGN.RIGHT)
        write(tf, "{:,}".format(total), size=15, bold=True, color=INK, font=F_MONO)
        write(tf, "answer @ {:,}".format(ans), size=9.5, color=INK2, font=F_MONO,
              para=tf.add_paragraph(), align=PP_ALIGN.RIGHT)
    ay = y + len(rows) * (row_h + gap) - gap + 0.10
    tf = textbox(s, track_x, ay, 2.0, 0.22)
    write(tf, "token 0", size=8.5, color=MUTE, font=F_MONO, spacing=0.8)
    tf = textbox(s, track_x + track_w - 2.0, ay, 2.0, 0.22, align=PP_ALIGN.RIGHT)
    write(tf, "token 7,460", size=8.5, color=MUTE, font=F_MONO, spacing=0.8)
    return ay + 0.4


def build(path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SLIDE_W), Inches(SLIDE_H)
    T = 12

    # 1 ── Title
    s = new_slide(prs, 1, T)
    tf = textbox(s, M, 2.15, CW, 0.3)
    write(tf, "Akamai · Edge content negotiation", size=10, bold=True,
          color=EDGE, font=F_MONO, spacing=1.8, caps=True)
    tf = textbox(s, M, 2.62, CW * 0.82, 1.7)
    write(tf, "Prioritized delivery\nfor AI answer engines", size=44, bold=True,
          color=INK, font=F_DISPLAY, line=1.02)
    tf = textbox(s, M, 4.5, CW * 0.62, 1.0)
    write(tf, "Your support content is being read by machines that answer on your "
              "behalf. This is how we control what they receive, in what order, and "
              "what it costs them to reach the answer — without changing your site.",
          size=13, color=INK2, font=F_BODY, line=1.5)
    for i, (tag, val, acc) in enumerate([
            ("Prepared for", "→ client name", EDGE),
            ("Track", "Prioritized delivery", None),
            ("Status", "Working mechanic", None)]):
        card(s, M + i * 2.75, 5.75, 2.5, 0.95, tag=tag, title=val, accent=acc)

    # 2 ── The shift
    s = new_slide(prs, 2, T)
    b = head(s, "The shift", "Three audiences now read your support pages.")
    items = [("Audience 01", "People",
              "Arrive in a browser, get the full experience. Unchanged by everything "
              "in this deck.", None),
             ("Audience 02", "Search crawlers",
              "Index your pages so they rank. Well understood, well served, and not "
              "the problem here.", None),
             ("Audience 03", "AI answer engines",
              "Read your pages to answer a customer's question instead of sending "
              "them to you. What they extract becomes your brand's answer — "
              "accurate or not.", EDGE)]
    for i, (tag, title, body, acc) in enumerate(items):
        card(s, M + i * 4.03, b + 0.15, 3.78, 2.3, tag=tag, title=title,
             body=body, accent=acc)
    tf = textbox(s, M, b + 2.75, CW * 0.78, 0.6)
    write(tf, "For the first two, the page is a destination. For the third, the page "
              "is source material. That changes what “well delivered” means.",
          size=13, color=INK2, font=F_BODY, line=1.5)

    # 3 ── The problem, measured
    s = new_slide(prs, 3, T)
    b = head(s, "The problem", "The answer is buried under the page.",
             lede="A representative support article — “how do I take a "
                  "screenshot” — measured end to end. A model reads top-down "
                  "and pays for every token on the way.")
    ry = ruler(s, b + 0.30, [("As published", 1.0, 0.664, 7460, 4955, BURIED)])
    stats = [("Before the answer", "4,955", BURIED,
              "tokens of cookie banner, mega-menu, promos, topic tiles, breadcrumbs, "
              "a marketing intro and a table of contents."),
             ("The answer itself", "4", INK,
              "steps. That is the entire thing the customer asked for."),
             ("Consequence", "66%", INK,
              "of the page is consumed before the model reaches anything worth citing.")]
    for i, (tag, val, col, body) in enumerate(stats):
        x = M + i * 4.03
        card(s, x, ry + 0.15, 3.78, 1.75, tag=tag)
        tf = textbox(s, x + 0.22, ry + 0.55, 3.3, 0.55)
        write(tf, val, size=30, bold=True, color=col, font=F_DISPLAY)
        tf = textbox(s, x + 0.22, ry + 1.12, 3.34, 0.7)
        write(tf, body, size=10, color=INK2, font=F_BODY, line=1.35)

    # 4 ── Why it matters
    s = new_slide(prs, 4, T)
    b = head(s, "Why it matters to marketing", "Buried answers get summarised, not cited.")
    card(s, M, b + 0.15, 7.1, 3.5, bullets=[
        "Context is finite. Every model has a budget. Pages that spend it on "
        "navigation get truncated, and the truncated part is often your actual answer.",
        "Position is weight. What appears first is treated as the point of the page. "
        "Your intro is not your point.",
        "Wrong beats absent, badly. A model that half-reads a device support page "
        "will still answer — confidently, and sometimes about the wrong handset.",
        "You already proved the content thesis. The open question was never whether "
        "better-structured content works. It was how to deliver it at scale."])
    x = M + 7.4
    card(s, x, b + 0.15, 4.43, 3.5, tag="Your PoC result", accent=REACHED)
    tf = textbox(s, x + 0.22, b + 1.05, 4.0, 0.6)
    write(tf, "→ 25–30% to 60%", size=25, bold=True, color=REACHED, font=F_DISPLAY)
    tf = textbox(s, x + 0.22, b + 1.75, 4.0, 1.5)
    write(tf, "citation accuracy, from restructuring content alone — hand-built, "
              "on a limited set of pages.\n\nEverything that follows is about doing "
              "that across your estate, automatically, without a two-month content "
              "project per template.", size=10.5, color=INK2, font=F_BODY, line=1.4)

    # 5 ── Routing table
    s = new_slide(prs, 5, T)
    b = head(s, "The mechanic", "One URL. The edge decides\nwhat each visitor receives.",
             lede="The request arrives at Akamai before it reaches you. We identify "
                  "the caller, then serve the representation that suits them — "
                  "from the same source content.")
    rows = [("A person", "no bot signal", "your site, untouched",
             "Nothing changes. No risk to the customer experience."),
            ("Googlebot, Bingbot", "verified bot signal", "full rendered HTML",
             "Search indexers want the crawlable page and its structured data."),
            ("ClaudeBot, GPTBot,\nGoogle-Extended", "verified bot signal", "skinny",
             "Ingesting for training. Chrome removed; original order kept."),
            ("PerplexityBot,\nOAI-SearchBot", "verified bot signal", "prioritized",
             "Answering a question right now. Answer hoisted to the top.")]
    tbl_h = 2.6
    gt = s.shapes.add_table(len(rows) + 1, 4, Inches(M), Inches(b + 0.15),
                            Inches(CW), Inches(tbl_h)).table
    gt.columns[0].width = Inches(2.7); gt.columns[1].width = Inches(2.2)
    gt.columns[2].width = Inches(2.5); gt.columns[3].width = Inches(4.43)
    for j, h in enumerate(["Who is asking", "Identified by", "Receives", "Why"]):
        c = gt.cell(0, j); c.text = ""
        tf = c.text_frame; tf.word_wrap = True
        write(tf, h, size=8.5, bold=True, color=MUTE, font=F_MONO, spacing=1.2, caps=True)
        c.fill.solid(); c.fill.fore_color.rgb = GROUND
    for i, r in enumerate(rows, start=1):
        for j, val in enumerate(r):
            c = gt.cell(i, j); c.text = ""
            tf = c.text_frame; tf.word_wrap = True
            write(tf, val, size=10,
                  color=INK if j == 0 else INK2,
                  font=F_MONO if j <= 2 else F_BODY,
                  bold=(j == 2), line=1.25)
            c.fill.solid(); c.fill.fore_color.rgb = PANEL
    note(s, "Identity comes from Akamai Bot Manager's verified bot signal, not the "
            "User-Agent string alone — a UA on its own is trivially spoofed. The "
            "policy is configurable per property; this table is the default.",
         b + tbl_h + 0.35)

    # 6 ── The result
    s = new_slide(prs, 6, T)
    b = head(s, "The result",
             "Stripping the page is the obvious win.\nReordering it is the free one.")
    ry = ruler(s, b + 0.25, [
        ("As published", 1.0,   0.664,  7460, 4955, BURIED),
        ("Skinny",       0.129, 0.0598,  961,  446, EDGE),
        ("Prioritized",  0.129, 0.0102,  960,   76, REACHED)])
    cards = [("Skinny", "Chrome, scripts, styles and media removed. Content order "
                        "untouched. −87% tokens.", None),
             ("Prioritized", "Everything skinny does, then the answer block is hoisted "
                             "above the intro and contents list. Answer 65× closer.", None),
             ("The point", "Skinny and prioritized are the same payload size — "
                           "3,953 bytes each. Ordering costs nothing and moves the "
                           "answer from token 446 to 76.", REACHED)]
    for i, (tag, body, acc) in enumerate(cards):
        card(s, M + i * 4.03, ry + 0.10, 3.78, 1.35, tag=tag, body=body, accent=acc)
    note(s, "Measured on a representative support-page structure, not yet on your live "
            "templates. Establishing your own baseline is step one of the rollout.",
         ry + 1.60)

    # 7 ── Architecture
    s = new_slide(prs, 7, T)
    head(s, "Architecture", "Request path, build path.", size=27)

    def band(label, y):
        tf = textbox(s, M, y, CW, 0.24)
        write(tf, label, size=8.5, color=MUTE, font=F_MONO, spacing=1.4, caps=True)

    def dbox(x, y, w, h, title, lines, accent=None):
        rect(s, x, y, w, h, fill=PANEL, border=accent or LINE,
             width=1.75 if accent else 1.25, rounded=True)
        tf = textbox(s, x + 0.18, y + 0.16, w - 0.36, 0.26)
        write(tf, title, size=11.5, bold=True, color=INK, font=F_DISPLAY)
        if lines:
            tf = textbox(s, x + 0.18, y + 0.46, w - 0.36, h - 0.6)
            for k, ln in enumerate(lines):
                write(tf, ln, size=9, color=INK2, font=F_BODY, line=1.3,
                      para=None if k == 0 else tf.add_paragraph())

    def alabel(text, x, y):
        tf = textbox(s, x, y, 1.0, 0.2)
        write(tf, text, size=8, color=MUTE, font=F_MONO, spacing=1.0, caps=True)

    band("Request path — milliseconds, in the hot path", 1.92)
    dbox(M, 2.22, 2.30, 1.20, "Answer engine",
         ["Perplexity, OAI-SearchBot,", "ClaudeBot, GPTBot…"])
    arrow(s, 3.10, 2.82, 3.68, 2.82, color=EDGE, width=2);  alabel("HTTPS", 3.14, 2.60)
    dbox(3.75, 2.12, 3.10, 1.40, "Akamai edge",
         ["1 · Verify bot identity", "2 · Apply routing policy",
          "3 · Select variant"], accent=EDGE)
    arrow(s, 6.92, 2.82, 7.48, 2.82, color=EDGE, width=2);  alabel("HIT", 7.00, 2.60)
    dbox(7.55, 2.12, 2.60, 1.40, "Variant cache",
         ["full · skinny · prioritized", "keyed by url + device",
          "+ lang + variant"], accent=REACHED)
    arrow(s, 10.22, 2.82, 10.78, 2.82, dash=True);          alabel("MISS", 10.26, 2.60)
    dbox(10.85, 2.32, 1.73, 1.00, "Your origin", ["unchanged"])

    band("Build path — ahead of time, never in the hot path", 3.98)
    dbox(M, 4.28, 2.30, 1.20, "Source content", ["your CMS / origin", "templates"])
    arrow(s, 3.10, 4.88, 3.68, 4.88)
    dbox(3.75, 4.14, 3.10, 1.55, "Variant builder",
         ["A · You supply the structured", "     version — we cache + serve it",
          "B · We derive it from the page", "     — rules, or model-generated"])
    arrow(s, 6.92, 4.88, 7.48, 4.88, color=REACHED);        alabel("WRITE", 6.96, 4.66)
    dbox(7.55, 4.28, 2.60, 1.20, "Variant cache", ["one build,", "served indefinitely"],
         accent=REACHED)
    arrow(s, 8.85, 4.28, 8.85, 3.56, color=REACHED)
    arrow(s, 10.78, 4.88, 10.22, 4.88, dash=True)
    dbox(10.85, 4.28, 1.73, 1.20, "50 configs", ["cover 375k", "pages"])

    note(s, "The two halves are deliberately separate. Building a variant is slow and "
            "can involve a model; serving one is a cache read. No generation ever "
            "happens while a crawler waits. Keeping the cached variant in step with "
            "the source is the third concern — that is the next slide.", 5.88)

    # 8 ── Producing the variants
    s = new_slide(prs, 8, T)
    b = head(s, "Producing the variants", "Two ways in. You are not committed to either.")
    card(s, M, b + 0.15, 5.79, 2.9, tag="Path A — you supply it",
         title="Bring your own grounding layer",
         body="You already know how to structure content for retrieval — that is "
              "what your PoC proved. Your team or your agent produces the structured "
              "version; we ingest, cache, route and keep it fresh.\n\n"
              "You keep editorial control. We are the delivery layer.")
    card(s, M + 6.04, b + 0.15, 5.79, 2.9, tag="Path B — we derive it",
         title="Generated from the page you already have",
         body="Rules-based per template — roughly 50 configurations covering the "
              "estate — or model-generated for the long tail where no template "
              "fits.\n\nNo content project required to get a baseline live.")
    tf = textbox(s, M, b + 3.30, CW * 0.78, 0.6)
    write(tf, "Most estates end up mixed: rules for the high-traffic templates, "
              "generation for the tail, hand-authored for the pages that matter most "
              "commercially.", size=13, color=INK2, font=F_BODY, line=1.5)

    # 9 ── Freshness
    s = new_slide(prs, 9, T)
    b = head(s, "The hard part",
             "A cached variant must never outlive\nthe page it came from.",
             lede="This is the objection worth spending time on, because a stale answer "
                  "about a tariff or a device is worse than no answer at all.")
    fr = [("Event-driven", "Purge on publish",
           "Your CMS already emits a publish event. It invalidates the variant at the "
           "same moment it invalidates the page. Same trigger, same instant."),
          ("Belt and braces", "Hash revalidation",
           "A background pass digests the source and compares. If the page changed "
           "without an event firing, the variant rebuilds itself."),
          ("Never cached", "Live holes",
           "Volatile fields — price, stock, plan availability — are marked "
           "and fetched at request time. They are never baked into a variant.")]
    for i, (tag, title, body) in enumerate(fr):
        card(s, M + i * 4.03, b + 0.20, 3.78, 2.2, tag=tag, title=title, body=body)
    note(s, "The failure mode we design against is not “the cache is slightly "
            "behind.” It is “an answer engine confidently quotes a tariff "
            "you withdrew last quarter.”", b + 2.60)

    # 10 ── Measurement
    s = new_slide(prs, 10, T)
    b = head(s, "What you get to measure", "Bot traffic stops being a blind spot.",
             lede="The edge sees every one of these requests. That makes a category of "
                  "reporting possible that most teams currently do not have at all.")
    ms = [("Who", "By engine", "Which answer engines crawl you, how often, and which "
                               "sections they favour."),
          ("What", "By variant", "What each one was served, and the cache hit rate "
                                 "behind it."),
          ("Cost", "By payload", "Tokens and bytes delivered — the real unit of "
                                 "what you are handing over."),
          ("Health", "By status", "Errors, staleness, and any page whose variant "
                                  "failed to build.")]
    for i, (tag, title, body) in enumerate(ms):
        card(s, M + i * 3.02, b + 0.20, 2.79, 1.95, tag=tag, title=title, body=body)
    tf = textbox(s, M, b + 2.40, CW * 0.8, 0.6)
    write(tf, "For a marketing team this is the feedback loop: change the structure, "
              "watch what the engines take, correlate with how your brand is being "
              "described.", size=13, color=INK2, font=F_BODY, line=1.5)

    # 11 ── Proven vs designed
    s = new_slide(prs, 11, T)
    b = head(s, "Where we actually are", "What is working today, and what is design.")
    card(s, M, b + 0.20, 5.79, 3.0, tag="Working, demonstrable now", accent=REACHED,
         bullets=["Bot identification and per-engine routing at the edge",
                  "Both transformations, measured — the numbers in this deck are live output",
                  "Model-generated variants for pages with no template config",
                  "Variant caching, and serving the right one per request"])
    card(s, M + 6.04, b + 0.20, 5.79, 3.0, tag="Designed, not yet built for you",
         bullets=["CMS purge integration against your publishing pipeline",
                  "Hash revalidation at your estate's scale",
                  "The reporting surface described on the previous slide",
                  "Template configs derived from your actual page structures"])
    note(s, "We would rather be precise about this line than blur it. The mechanic is "
            "real; the integration is scoped work.", b + 3.35)

    # 12 ── Rollout
    s = new_slide(prs, 12, T)
    b = head(s, "Proposed next step",
             "Prove it on one section before\nit touches anything else.")
    steps = [("Step 01 · weeks 1–2", "Baseline",
              "Measure your real templates as published: tokens, answer position, what "
              "the engines currently receive. No changes yet."),
             ("Step 02 · weeks 3–6", "One section, live",
              "Device support, or whichever section carries the most answerable "
              "questions. Routing on, variants cached, freshness wired to your "
              "publish events."),
             ("Step 03 · ongoing", "Measure, then widen",
              "Compare citation behaviour against the untouched sections. Expand "
              "template by template on evidence.")]
    for i, (tag, title, body) in enumerate(steps):
        card(s, M + i * 4.03, b + 0.20, 3.78, 2.0, tag=tag, title=title, body=body)
    card(s, M, b + 2.35, 5.79, 1.25, tag="Why a section, not the estate", accent=EDGE,
         body="It bounds the blast radius, it gives you a clean control group, and it "
              "produces a number you can defend internally before anyone commits "
              "budget to the rest.")
    card(s, M + 6.04, b + 2.35, 5.79, 1.25, tag="What we need from you",
         body="A section to target, read access to the templates, and whoever owns the "
              "publish pipeline for thirty minutes.")

    prs.save(path)
    return path


if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else "prioritized-delivery.pptx"
    print("wrote", build(out))
